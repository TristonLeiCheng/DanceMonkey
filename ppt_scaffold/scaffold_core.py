"""
DanceMonkey 内置 python-pptx 脚手架：多主题、浅色底、顶栏强调、内容卡片。
与 Agent 系统提示中的「禁止纯黑默认底」一致；请在此基础上改主题名与文案。
"""
from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class Theme:
    name: str
    slide_bg: RGBColor
    card_bg: RGBColor
    accent: RGBColor
    accent2: RGBColor
    text_primary: RGBColor
    text_secondary: RGBColor
    text_muted: RGBColor
    rule: RGBColor


THEMES: dict[str, Theme] = {
    "light_business": Theme(
        name="light_business",
        slide_bg=RGBColor(0xF4, 0xF6, 0xFA),
        card_bg=RGBColor(0xFF, 0xFF, 0xFF),
        accent=RGBColor(0x3B, 0x5B, 0xDB),
        accent2=RGBColor(0x14, 0xB8, 0xA6),
        text_primary=RGBColor(0x11, 0x18, 0x27),
        text_secondary=RGBColor(0x37, 0x41, 0x51),
        text_muted=RGBColor(0x6B, 0x72, 0x80),
        rule=RGBColor(0xE2, 0xE8, 0xF0),
    ),
    "tech_blue": Theme(
        name="tech_blue",
        slide_bg=RGBColor(0xEC, 0xF4, 0xFF),
        card_bg=RGBColor(0xFF, 0xFF, 0xFF),
        accent=RGBColor(0x25, 0x63, 0xEB),
        accent2=RGBColor(0x22, 0xD3, 0xEE),
        text_primary=RGBColor(0x0F, 0x17, 0x2A),
        text_secondary=RGBColor(0x1E, 0x2A, 0x3B),
        text_muted=RGBColor(0x64, 0x74, 0x8B),
        rule=RGBColor(0xBF, 0xDB, 0xFE),
    ),
    "warm_magazine": Theme(
        name="warm_magazine",
        slide_bg=RGBColor(0xFF, 0xF7, 0xED),
        card_bg=RGBColor(0xFF, 0xFF, 0xFF),
        accent=RGBColor(0xEA, 0x58, 0x0C),
        accent2=RGBColor(0xD9, 0x77, 0x06),
        text_primary=RGBColor(0x29, 0x1C, 0x0F),
        text_secondary=RGBColor(0x44, 0x2C, 0x1A),
        text_muted=RGBColor(0x78, 0x56, 0x3A),
        rule=RGBColor(0xFE, 0xD7, 0xAA),
    ),
}


def new_presentation_16_9() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_top_bar(slide, theme: Theme, height_in: float = 0.12) -> None:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs_slide_width(slide),
        Inches(height_in),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = theme.accent
    sh.line.fill.background()


def prs_slide_width(slide) -> int:
    try:
        return slide.part.presentation_part.presentation.slide_width
    except AttributeError:
        return int(Inches(13.333))


def add_footer(slide, theme: Theme, left_text: str, right_text: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = left_text
    p.font.size = Pt(9)
    p.font.color.rgb = theme.text_muted
    p.alignment = PP_ALIGN.LEFT
    if right_text:
        p2 = tf.add_paragraph()
        p2.text = right_text
        p2.font.size = Pt(9)
        p2.font.color.rgb = theme.text_muted
        p2.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation, theme_name: str, kicker: str, title: str, subtitle: str) -> None:
    theme = THEMES[theme_name]
    slide = blank_slide(prs)
    set_slide_bg(slide, theme.slide_bg)
    add_top_bar(slide, theme)

    kb = slide.shapes.add_textbox(Inches(0.75), Inches(0.95), Inches(11.8), Inches(0.35))
    p0 = kb.text_frame.paragraphs[0]
    p0.text = kicker.upper()
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = theme.text_muted
    p0.alignment = PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(11.5), Inches(1.2))
    tfp = tb.text_frame
    tfp.word_wrap = True
    p1 = tfp.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = theme.text_primary
    p1.line_spacing = 1.05

    sb = slide.shapes.add_textbox(Inches(0.75), Inches(2.85), Inches(11.2), Inches(0.9))
    p2 = sb.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = theme.text_secondary

    # 右侧装饰卡
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.35), Inches(3.6), Inches(4.2))
    card.fill.solid()
    card.fill.fore_color.rgb = theme.card_bg
    card.line.color.rgb = theme.rule

    deco = slide.shapes.add_textbox(Inches(9.45), Inches(2.0), Inches(3.1), Inches(2.8))
    dp = deco.text_frame.paragraphs[0]
    dp.text = "DanceMonkey\nDeck Scaffold"
    dp.font.size = Pt(14)
    dp.font.bold = True
    dp.font.color.rgb = theme.accent
    dp.alignment = PP_ALIGN.CENTER
    deco.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_footer(slide, theme, "Built-in scaffold · " + theme.name, "")


def add_content_slide(
    prs: Presentation,
    theme_name: str,
    index_text: str,
    title: str,
    bullets: Sequence[str],
) -> None:
    theme = THEMES[theme_name]
    slide = blank_slide(prs)
    set_slide_bg(slide, theme.slide_bg)
    add_top_bar(slide, theme, 0.08)

    idx = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(1.2), Inches(0.4))
    ip = idx.text_frame.paragraphs[0]
    ip.text = index_text
    ip.font.size = Pt(12)
    ip.font.bold = True
    ip.font.color.rgb = theme.accent2

    tt = slide.shapes.add_textbox(Inches(0.65), Inches(1.0), Inches(11.8), Inches(0.75))
    tp = tt.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = theme.text_primary

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.82), Inches(12.0), Inches(0.02))
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.rule
    rule.line.fill.background()

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.05), Inches(12.0), Inches(4.85))
    card.fill.solid()
    card.fill.fore_color.rgb = theme.card_bg
    card.line.color.rgb = theme.rule

    body = slide.shapes.add_textbox(Inches(0.95), Inches(2.35), Inches(11.4), Inches(4.35))
    tf = body.text_frame
    tf.word_wrap = True
    if bullets:
        p0 = tf.paragraphs[0]
        p0.text = bullets[0].strip()
        p0.level = 0
        p0.font.size = Pt(15)
        p0.font.color.rgb = theme.text_secondary
        p0.space_after = Pt(6)
        p0.bullet = True
        for line in bullets[1:]:
            pn = tf.add_paragraph()
            pn.text = line.strip()
            pn.level = 0
            pn.font.size = Pt(15)
            pn.font.color.rgb = theme.text_secondary
            pn.space_after = Pt(6)
            pn.bullet = True

    add_footer(slide, theme, "DanceMonkey", index_text)


def add_two_column_slide(
    prs: Presentation,
    theme_name: str,
    index_text: str,
    title: str,
    left_bullets: Sequence[str],
    right_bullets: Sequence[str],
    left_header: str = "",
    right_header: str = "",
) -> None:
    """双栏版式：左右两组要点并列，适合对比/方案/优缺点等场景。"""
    theme = THEMES[theme_name]
    slide = blank_slide(prs)
    set_slide_bg(slide, theme.slide_bg)
    add_top_bar(slide, theme, 0.08)

    # 标题区
    idx = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(1.2), Inches(0.4))
    ip = idx.text_frame.paragraphs[0]
    ip.text = index_text
    ip.font.size = Pt(12)
    ip.font.bold = True
    ip.font.color.rgb = theme.accent2

    tt = slide.shapes.add_textbox(Inches(0.65), Inches(1.0), Inches(11.8), Inches(0.75))
    tp = tt.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = theme.text_primary

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.82), Inches(12.0), Inches(0.02))
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.rule
    rule.line.fill.background()

    # 左栏 accent 竖条
    lbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.05), Inches(0.05), Inches(4.85))
    lbar.fill.solid()
    lbar.fill.fore_color.rgb = theme.accent
    lbar.line.fill.background()

    # 右栏 accent 竖条
    rbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.05), Inches(0.05), Inches(4.85))
    rbar.fill.solid()
    rbar.fill.fore_color.rgb = theme.accent2
    rbar.line.fill.background()

    def _build_column(x: float, header: str, items: Sequence[str]) -> None:
        if header:
            hb = slide.shapes.add_textbox(Inches(x), Inches(2.05), Inches(5.7), Inches(0.4))
            hp = hb.text_frame.paragraphs[0]
            hp.text = header
            hp.font.size = Pt(14)
            hp.font.bold = True
            hp.font.color.rgb = theme.accent
            body_top = 2.55
        else:
            body_top = 2.2

        if not items:
            return
        body = slide.shapes.add_textbox(Inches(x), Inches(body_top), Inches(5.7), Inches(4.6 - (body_top - 2.2)))
        tf = body.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = items[0].strip()
        p0.font.size = Pt(14)
        p0.font.color.rgb = theme.text_secondary
        p0.space_after = Pt(5)
        p0.bullet = True
        for line in items[1:]:
            pn = tf.add_paragraph()
            pn.text = line.strip()
            pn.font.size = Pt(14)
            pn.font.color.rgb = theme.text_secondary
            pn.space_after = Pt(5)
            pn.bullet = True

    _build_column(0.85, left_header, left_bullets)
    _build_column(7.0, right_header, right_bullets)

    add_footer(slide, theme, "DanceMonkey", index_text)


def add_quote_slide(
    prs: Presentation,
    theme_name: str,
    quote_text: str,
    source: str = "",
    page_title: str = "",
) -> None:
    """引述版式：大字号金句/名言/数据亮点，适合强调单一核心信息。"""
    theme = THEMES[theme_name]
    slide = blank_slide(prs)
    set_slide_bg(slide, theme.slide_bg)
    add_top_bar(slide, theme)

    # 左侧粗 accent 竖条（视觉锚点）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(0.1), Inches(4.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent
    bar.line.fill.background()

    # 大引号装饰
    qq = slide.shapes.add_textbox(Inches(0.85), Inches(1.1), Inches(1.0), Inches(0.85))
    qp = qq.text_frame.paragraphs[0]
    qp.text = "\u201C"
    qp.font.size = Pt(64)
    qp.font.bold = True
    qp.font.color.rgb = theme.accent2

    # 引述正文
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.2), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote_text
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = theme.text_primary
    p.line_spacing = 1.25

    # 出处
    if source:
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.9), Inches(2.2), Inches(0.02))
        rule.fill.solid()
        rule.fill.fore_color.rgb = theme.rule
        rule.line.fill.background()

        sb = slide.shapes.add_textbox(Inches(0.9), Inches(5.05), Inches(8.0), Inches(0.45))
        sp = sb.text_frame.paragraphs[0]
        sp.text = f"— {source}"
        sp.font.size = Pt(13)
        sp.font.color.rgb = theme.text_muted

    # 可选标题（小字，右下方）
    if page_title:
        ptb = slide.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(10.0), Inches(0.4))
        ptp = ptb.text_frame.paragraphs[0]
        ptp.text = page_title
        ptp.font.size = Pt(11)
        ptp.font.bold = True
        ptp.font.color.rgb = theme.text_muted

    add_footer(slide, theme, "DanceMonkey", "")


def add_closing_slide(prs: Presentation, theme_name: str, headline: str, subline: str) -> None:
    theme = THEMES[theme_name]
    slide = blank_slide(prs)
    set_slide_bg(slide, theme.slide_bg)
    add_top_bar(slide, theme)

    t = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    p.text = headline
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme.text_primary
    p.alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.6))
    sp = s.text_frame.paragraphs[0]
    sp.text = subline
    sp.font.size = Pt(18)
    sp.font.color.rgb = theme.text_secondary
    sp.alignment = PP_ALIGN.CENTER

    add_footer(slide, theme, "Thank you", "")
